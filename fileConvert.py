from flask import Flask, request, jsonify, render_template, redirect, url_for
import os
import pandas as pd
from PIL import Image
import pytesseract
import PyPDF2
import json
import numpy as np
import logging
import io
import hashlib
import time
import random
import re
import requests

app = Flask(__name__)
UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Store combined data temporarily
global combined_data_global
combined_data_global = []

# Define path to Tesseract executable
pytesseract.pytesseract.tesseract_cmd = '/usr/bin/tesseract'

# Function to process different file types
def extract_data(file_path, file_type):
    if file_type == 'csv':
        # Use pandas to read csv file
        with open(file_path, 'r') as file:
            lines = file.readlines()
        addresses = [line.strip() for line in lines[1:] if line.strip()]  # Skip the header and strip whitespace
        return addresses

    elif file_type == 'xlsx':
        # Use pandas to read Excel file
        df = pd.read_excel(file_path)
        df = df.replace({np.nan: None})  # Replace NaN with None to make it JSON serializable
        return df.to_dict(orient='records')

    elif file_type == 'pdf':
        # Use PyPDF2 to read pdf file
        pdf_data = []
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                pdf_data.append(page.extract_text())
        # Join all pages into a single text block
        full_text = " ".join(pdf_data)
        return full_text

    elif file_type in ['jpeg', 'jpg', 'png']:
        # Use pytesseract to extract text from image
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
        return text

    elif file_type in ['ppt', 'pptx']:
        import pptx
        ppt_data = []
        prs = pptx.Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    ppt_data.append(shape.text)
        # Join all slides into a single text block
        full_text = " ".join([re.sub(r'\s+', ' ', slide.strip()) for slide in ppt_data if slide.strip()])
        return full_text

    else:
        return {'error': 'Unsupported file type'}

# Route for uploading and processing the file
@app.route('/')
def index():
    return '''
    <!DOCTYPE html>
    <html lang="en">
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>Upload File</title>
    </head>
    <body>
        <h1>Upload a File</h1>
        <form action="/upload" method="post" enctype="multipart/form-data">
            <input type="file" name="file" required />
            <input type="submit" value="Upload" />
        </form>
    </body>
    </html>
    '''

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    # Save the file to a folder and process
    if file:
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
        file.save(file_path)

        # Extract data based on file type
        file_type = file.filename.rsplit('.', 1)[1].lower()
        raw_data = extract_data(file_path, file_type)

        # Log the JSON data
        logger.info(f"JSON data extracted: {raw_data}")

        # Send data to webhook
        webhook_url = "https://hook.us1.make.com/huolkx7l5lpug0q51wxftsvfctnkcday"
        payload = {'data': raw_data}
        headers = {
            'Content-Type': 'application/json'
        }
        try:
            response = requests.post(webhook_url, json=payload, headers=headers)
            response.raise_for_status()
            logger.info(f"Data successfully sent to webhook: {response.status_code}")
        except requests.exceptions.RequestException as e:
            logger.error(f"Failed to send data to webhook: {e}")
            return jsonify({'error': 'Failed to send data to webhook'}), 400

        # Store the data in global variable for displaying in table later
        global combined_data_global
        combined_data_global = []  # Reset to ensure we start fresh

        # Redirect to the /loading route
        return redirect(url_for('loading'))

@app.route('/loading')
def loading():
    return '''
    <!DOCTYPE html>
    <html lang="en">
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>Processing...</title>
        <script>
            function checkDataReady() {
                fetch('/is_data_ready')
                    .then(response => response.json())
                    .then(data => {
                        if (data.ready) {
                            window.location.href = "/webhook";
                        } else {
                            setTimeout(checkDataReady, 1000); // Check again after 1 second
                        }
                    })
                    .catch(error => {
                        console.error('Error checking data readiness:', error);
                    });
            }
            
            // Start checking if data is ready
            checkDataReady();
        </script>
    </head>
    <body>
        <h2>Please wait while your data is being processed...</h2>
        <div id="loading">Loading...</div>
    </body>
    </html>
    '''

@app.route('/is_data_ready')
def is_data_ready():
    global combined_data_global
    if combined_data_global:
        return jsonify({'ready': True})
    return jsonify({'ready': False})

# Route to act as a webhook to receive data and generate a table
@app.route('/webhook', methods=['POST', 'GET'])
def display_table():
    global combined_data_global
    try:
        # For GET requests, use the previously saved data
        if request.method == 'GET':
            if not combined_data_global:
                return '''
                <!DOCTYPE html>
                <html lang="en">
                <head>
                    <meta charset="UTF-8">
                    <meta name="viewport" content="width=device-width, initial-scale=1.0">
                    <title>Property Estimates Table</title>
                    <style>
                        a.convert-new {{
                            position: absolute;
                            top: 20px;
                            right: 20px;
                            text-decoration: none;
                            padding: 10px;
                            background-color: #007bff;
                            color: white;
                            border-radius: 5px;
                        }}
                    </style>
                </head>
                <body>
                    <a href="/" class="convert-new">Convert new file</a>
                    <h2>Property Estimates</h2>
                    <p>No data available to display.</p>
                </body>
                </html>
                '''
            incoming_data = combined_data_global
        else:
            incoming_data = request.get_json()
            if not incoming_data:
                return jsonify({'error': 'No data received'}), 400

        # Handle incoming data if it's a list
        if isinstance(incoming_data, list):
            combined_data = incoming_data
        else:
            logger.info(f"Received data from webhook: {incoming_data}")

            # Extract ChatGPT data and estimates data
            chatgpt_data = incoming_data.get("chatgpt_data", {}).get("properties", [])
            estimate_data = incoming_data.get("estimate_data", [])

            # Merge the data sets based on "full_address"
            combined_data = []

            # Convert estimate_data to a dictionary for easier lookups
            estimates_dict = {item['full_address']: item for item in estimate_data}

            for item in chatgpt_data:
                full_address = item.get("full_address")
                combined_entry = item.copy()  # Start with chatgpt_data
                if full_address and full_address in estimates_dict:
                    # Add estimates to the entry
                    combined_entry.update(estimates_dict[full_address])
                combined_data.append(combined_entry)

            # Ensure that estimate data rows are appended at the end
            combined_data.extend([value for key, value in estimates_dict.items() if key not in [item.get("full_address") for item in chatgpt_data]])

        # Update the global variable with combined data
        combined_data_global = combined_data

        # Convert the combined data to a pandas DataFrame to create a table
        df = pd.DataFrame(combined_data)

        # Reformat column headers to readable text and make them sentence case
        df.columns = [re.sub(r'_', ' ', col).capitalize() for col in df.columns]

        # Replace NaN with empty strings in the dataframe
        df = df.replace({np.nan: ''})

        # Generate an HTML table representation with styling and filtering capabilities
        if not df.empty:
            html_table = df.to_html(index=False, classes='table table-bordered table-hover cell-border', border=0)
            html_table = f"<table id=\"propertyTable\">{html_table.split('>', 1)[1]}"  # Adding id for the table and removing unnecessary attributes
        else:
            html_table = "<p>No data available to display.</p>"

        # Log the combined data and return as response
        logger.info(f"Combined data: {combined_data}")
        return f'''
        <!DOCTYPE html>
        <html lang="en">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>Property Estimates Table</title>
            <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/datatables/1.10.21/css/jquery.dataTables.min.css">
            <script src="https://code.jquery.com/jquery-3.5.1.js"></script>
            <script src="https://cdn.datatables.net/1.10.21/js/jquery.dataTables.min.js"></script>
            <style>
                a.convert-new {{
                    position: absolute;
                    top: 20px;
                    right: 20px;
                    text-decoration: none;
                    padding: 10px;
                    background-color: #007bff;
                    color: white;
                    border-radius: 5px;
                }}
                th {{ text-align: center; }}
                table {{ width: 100%; }}
                table, th, td {{ border: 1px solid black; border-collapse: collapse; }}
            </style>
            <script>
                $(document).ready(function() {{
                    $('#propertyTable').DataTable({{
                        "paging": true,
                        "searching": true,
                        "ordering": true,
                    }});
                }});
            </script>
        </head>
        <body>
            <a href="/" class="convert-new">Convert new file</a>
            <h2>Property Estimates</h2>
            {html_table}
        </body>
        </html>
        '''

    except Exception as e:
        logger.error(f"Error receiving data: {e}")
        return jsonify({'error': 'Failed to receive data'}), 400

if __name__ == '__main__':
    port = int(os.environ.get("PORT", 5000))  # Use Heroku's PORT or default to 5000
    app.run(debug=False, host='0.0.0.0', port=port)
